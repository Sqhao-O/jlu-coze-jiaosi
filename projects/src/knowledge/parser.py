"""
文档解析模块

支持 PDF / Word / Excel / PPT / 纯文本 的文本提取。
"""

import logging
import os
from io import BytesIO
from typing import Optional

logger = logging.getLogger(__name__)


def extract_text(content: bytes, filename: str) -> str:
    """
    根据文件扩展名选择解析器，提取纯文本。

    Args:
        content: 文件二进制内容
        filename: 原始文件名（用于判断扩展名）

    Returns:
        提取的纯文本
    """
    ext = os.path.splitext(filename)[1].lower()

    try:
        if ext == ".pdf":
            return _extract_pdf(content)
        elif ext in (".docx", ".doc"):
            return _extract_docx(content)
        elif ext in (".xlsx", ".xls"):
            return _extract_xlsx(content)
        elif ext in (".pptx", ".ppt"):
            return _extract_pptx(content)
        elif ext == ".txt":
            return content.decode("utf-8", errors="replace")
        elif ext == ".md":
            return content.decode("utf-8", errors="replace")
        else:
            logger.warning(f"[文档解析] 不支持的文件类型: {ext}")
            return content.decode("utf-8", errors="replace")
    except Exception as e:
        logger.error(f"[文档解析] 解析 {filename} 失败: {e}")
        return ""


def _extract_pdf(content: bytes) -> str:
    """从 PDF 提取文本"""
    from pypdf import PdfReader

    reader = PdfReader(BytesIO(content))
    texts = []
    for page in reader.pages:
        text = page.extract_text()
        if text:
            texts.append(text.strip())
    return "\n\n".join(texts)


def _extract_docx(content: bytes) -> str:
    """从 Word 文档提取文本"""
    from docx2python import docx2python

    result = docx2python(BytesIO(content))
    return result.text


def _extract_xlsx(content: bytes) -> str:
    """从 Excel 提取文本（按 sheet + 行拼接）"""
    from openpyxl import load_workbook

    wb = load_workbook(BytesIO(content), read_only=True, data_only=True)
    texts = []
    for sheet in wb.worksheets:
        for row in sheet.iter_rows(values_only=True):
            cells = [str(c).strip() for c in row if c is not None]
            if cells:
                texts.append(" | ".join(cells))
    wb.close()
    return "\n".join(texts)


def _extract_pptx(content: bytes) -> str:
    """从 PPT 提取文本"""
    from pptx import Presentation

    prs = Presentation(BytesIO(content))
    texts = []
    for slide in prs.slides:
        slide_texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        slide_texts.append(text)
        if slide_texts:
            texts.append("\n".join(slide_texts))
    return "\n\n".join(texts)
